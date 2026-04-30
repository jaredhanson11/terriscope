"""Builds a .pptx territory report from slide records and S3 images.

Loads `ppt_assets/template.pptx` as a starting deck. The template has two slides:
  Slide 0 — Cover. Has a text shape named `title_slot` whose text is replaced
            with `cover_title` in place; styling is preserved from the template.
  Slide 1 — Content sample. Has named marker shapes `title_slot`, `image_slot`,
            `table_slot`. The slide is duplicated once per content slide. Each
            duplicate gets its title text replaced, then `image_slot` and
            `table_slot` are read for geometry, replaced with real content,
            and removed. The original template content slide is deleted at the
            end so it does not appear in the output.

Designers edit `template.pptx` in PowerPoint (colors, fonts, logo, marker
positions). The code only depends on shape *names*, so layout changes do not
require code changes.
"""

from collections.abc import Sequence
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import Any

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # type: ignore[attr-defined]
from pptx.util import Emu, Inches, Length, Pt

from src.models.exports import MapExportSlideModel
from src.services.s3 import S3Service

_TEMPLATE_PATH = Path(__file__).parent / "ppt_assets" / "template.pptx"

# Table styling
_HEADER_BG = RGBColor(0x6D, 0x4C, 0xE8)
_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
_ROW_BG = RGBColor(0xFF, 0xFF, 0xFF)
_ROW_BG_ALT = RGBColor(0xF8, 0xF9, 0xFB)
_BODY_FG = RGBColor(0x11, 0x18, 0x27)


def _find_shape(slide: Any, name: str) -> Any:
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None


def _replace_text(shape: Any, text: str) -> None:
    """Replace shape text in place, preserving the first run's font formatting."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for extra_run in list(p.runs[1:]):
            extra_run._r.getparent().remove(extra_run._r)
    else:
        p.text = text
    for extra_p in list(tf.paragraphs[1:]):
        extra_p._p.getparent().remove(extra_p._p)


def _remove_shape(shape: Any) -> None:
    shape.element.getparent().remove(shape.element)


def _duplicate_slide(prs: Any, src: Any) -> Any:
    """Append a deep copy of `src` to the deck and return it.

    python-pptx has no native duplicate; we add a fresh slide using the same
    layout (which gives us the correct rels), wipe the inherited placeholders,
    then deep-copy each shape's XML from the source.
    """
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shp in list(new_slide.shapes):
        _remove_shape(shp)
    for shp in src.shapes:
        new_slide.shapes._spTree.append(deepcopy(shp.element))
    return new_slide


def _delete_slide_at(prs: Any, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_id_lst.remove(list(sld_id_lst)[index])


def _fit_in_box(img_bytes: bytes, box_w: int, box_h: int) -> tuple[int, int, int, int]:
    """Return (left_offset, top_offset, width, height) in EMUs to center an image in a box."""
    with PILImage.open(BytesIO(img_bytes)) as img:
        iw, ih = img.size
    scale = min(box_w / iw, box_h / ih)
    w = round(iw * scale)
    h = round(ih * scale)
    return (box_w - w) // 2, (box_h - h) // 2, w, h


def _format_value(val: Any) -> str:
    if val is None:
        return ""
    if isinstance(val, bool):
        return "Yes" if val else "No"
    if isinstance(val, int):
        return f"{val:,}"
    if isinstance(val, float):
        return f"{val:,.2f}"
    return str(val)


def _is_numeric(val: Any) -> bool:
    return isinstance(val, (int, float)) and not isinstance(val, bool)


def _table_font_sizes(row_count: int, col_count: int) -> tuple[Any, Any]:
    """Pick header/body point sizes that scale down as rows/columns grow.

    Tables sit in a fixed slot, so wide configs (many columns) and tall configs
    (many rows) both need smaller text to keep cells readable without overflow.
    """
    if col_count >= 7 or row_count > 24:
        return Pt(8), Pt(7)
    if col_count >= 5 or row_count > 16:
        return Pt(9), Pt(8)
    return Pt(10), Pt(9)


def _build_table(
    slide: Any,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    node_data: list[dict[str, Any]],
) -> None:
    if not node_data:
        return

    # Column keys come from the first row's dict-insertion order. Callers in
    # ppt_exports.py emit display-ready keys ("Population (Sum)", etc.), so we
    # use them verbatim as headers — no snake_case-to-title transform.
    first = node_data[0]
    data_keys = [k for k in first if k != "name"]
    raw_keys = ["name", *data_keys]
    headers = ["Name", *data_keys]
    numeric_cols = {
        ci for ci, key in enumerate(raw_keys) if any(_is_numeric(n.get(key)) for n in node_data)
    }
    header_size, body_size = _table_font_sizes(len(node_data), len(headers))

    tbl = slide.shapes.add_table(
        rows=len(node_data) + 1,
        cols=len(headers),
        left=left,
        top=top,
        width=width,
        height=height,
    ).table

    for ci, header in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _HEADER_BG
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        para = cell.text_frame.paragraphs[0]
        para.text = header
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.bold = True
        run.font.size = header_size
        run.font.color.rgb = _HEADER_FG

    for ri, node in enumerate(node_data, start=1):
        zebra = ri % 2 == 0
        for ci, key in enumerate(raw_keys):
            val = node.get(key)
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ROW_BG_ALT if zebra else _ROW_BG
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            para = cell.text_frame.paragraphs[0]
            para.text = _format_value(val)
            para.alignment = PP_ALIGN.RIGHT if ci in numeric_cols and ci != 0 else PP_ALIGN.LEFT
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = body_size
            run.font.color.rgb = _BODY_FG


def _fill_content_slide(
    slide: Any,
    title: str,
    image_bytes: bytes | None,
) -> None:
    title_shape = _find_shape(slide, "title_slot")
    if title_shape is not None:
        _replace_text(title_shape, title)

    img_slot = _find_shape(slide, "image_slot")
    if img_slot is not None:
        if image_bytes:
            box_l, box_t, box_w, box_h = img_slot.left, img_slot.top, img_slot.width, img_slot.height
            off_l, off_t, w, h = _fit_in_box(image_bytes, int(box_w), int(box_h))
            slide.shapes.add_picture(
                BytesIO(image_bytes),
                left=Emu(int(box_l) + off_l),
                top=Emu(int(box_t) + off_t),
                width=Emu(w),
                height=Emu(h),
            )
        _remove_shape(img_slot)


def _fill_content_table(slide: Any, node_data: list[dict[str, Any]]) -> None:
    tbl_slot = _find_shape(slide, "table_slot")
    if tbl_slot is None:
        return
    left = Emu(int(tbl_slot.left))
    top = Emu(int(tbl_slot.top))
    width = Emu(int(tbl_slot.width))
    height = Emu(int(tbl_slot.height))
    _remove_shape(tbl_slot)
    _build_table(slide, left, top, width, height, node_data)


def build_pptx_buffer(
    slides: Sequence[MapExportSlideModel],
    s3: S3Service,
    cover_title: str = "Territory Report",
    layer_names: dict[int, str] | None = None,
) -> BytesIO:
    """Build the .pptx in memory and return a seek-zeroed buffer.

    `cover_title` populates the `title_slot` text on the cover slide.
    `layer_names` maps layer id → display name; when provided, a divider slide
    (cloned from `divider_template`) is inserted before the first content slide
    of each layer to separate the deck into sections.
    """
    prs = Presentation(str(_TEMPLATE_PATH))
    cover_slide = prs.slides[0]
    content_template = prs.slides[1]
    divider_template = prs.slides[2] if len(prs.slides) >= 3 else None

    cover_title_shape = _find_shape(cover_slide, "title_slot")
    if cover_title_shape is not None:
        _replace_text(cover_title_shape, cover_title)

    last_layer_id: int | None = None
    for slide_model in slides:
        if (
            divider_template is not None
            and layer_names is not None
            and slide_model.layer_id != last_layer_id
        ):
            section_title = layer_names.get(slide_model.layer_id, "")
            if section_title:
                divider = _duplicate_slide(prs, divider_template)
                divider_title_shape = _find_shape(divider, "title_slot")
                if divider_title_shape is not None:
                    _replace_text(divider_title_shape, section_title)
            last_layer_id = slide_model.layer_id

        image_bytes: bytes | None = None
        if slide_model.image_s3_key:
            image_bytes = s3.get_private_object(key=slide_model.image_s3_key).read()

        new_slide = _duplicate_slide(prs, content_template)
        _fill_content_slide(new_slide, title=slide_model.title or "", image_bytes=image_bytes)
        _fill_content_table(new_slide, slide_model.node_data or [])

    # Remove the original template source slides (content at index 1, divider
    # at index 2) — they were only used as duplication sources and should not
    # appear in the output deck. Delete from highest index to lowest so the
    # earlier index stays valid after the first removal.
    if divider_template is not None:
        _delete_slide_at(prs, 2)
    _delete_slide_at(prs, 1)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
