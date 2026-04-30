"""Generate the initial pptx export template.

Run once to produce src/services/ppt_assets/template.pptx. The result is
committed and edited by designers in PowerPoint from there on; this script is
just a starting point and a way to reset the template if it gets corrupted.

The template is a 3-slide deck:
  Slide 1 — Cover: a single named text shape `title_slot` whose text is
            replaced with the map name at export time.
  Slide 2 — Content sample: a `title_slot` text shape plus two named
            rectangle markers `image_slot` and `table_slot` whose geometry
            tells the builder where to drop the screenshot and the data
            table. The marker rectangles are removed at export time.
  Slide 3 — Divider sample: a single `title_slot` text shape used as a
            section separator before each layer's content slides
            (e.g., "Areas", "Regions", "Territories").

Usage:
    cd api && poetry run python scripts/build_pptx_template.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN  # type: ignore[attr-defined]
from pptx.util import Inches, Pt

# Slide size — keep at the existing 10 x 5.625 (16:9) so existing screenshots
# keep their aspect ratio. Designers can resize later if they want a bigger canvas.
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Brand colors — derived from the app's --primary token (oklch(0.51 0.23 277)).
PRIMARY = RGBColor(0x6D, 0x4C, 0xE8)
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
MARKER_FILL = RGBColor(0xF3, 0xF4, 0xF6)
MARKER_LINE = RGBColor(0xD1, 0xD5, 0xDB)


def _set_text(
    text_frame,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_brand_bar(slide, top, height) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY


def _add_marker(slide, name: str, left, top, width, height, label: str) -> None:
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.name = name
    sh.fill.solid()
    sh.fill.fore_color.rgb = MARKER_FILL
    sh.line.color.rgb = MARKER_LINE
    sh.line.width = Pt(0.75)
    _set_text(sh.text_frame, label, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def _build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_brand_bar(slide, 0, Inches(0.12))
    _add_brand_bar(slide, SLIDE_H - Inches(0.12), Inches(0.12))

    eyebrow = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), SLIDE_W - Inches(1.2), Inches(0.4))
    _set_text(eyebrow.text_frame, "TERRITORY REPORT", size=14, bold=True, color=PRIMARY)

    title = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), SLIDE_W - Inches(1.2), Inches(1.6))
    title.name = "title_slot"
    _set_text(title.text_frame, "{{cover_title}}", size=40, bold=True, color=INK)


def _build_content(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_bar(slide, 0, Inches(0.08))

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.18), SLIDE_W - Inches(0.6), Inches(0.5))
    title.name = "title_slot"
    _set_text(title.text_frame, "{{title}}", size=20, bold=True, color=INK)

    # Image slot — left ~2/3
    img_left = Inches(0.3)
    img_top = Inches(0.78)
    img_w = Inches(6.2)
    img_h = SLIDE_H - img_top - Inches(0.25)
    _add_marker(slide, "image_slot", img_left, img_top, img_w, img_h, "Map screenshot")

    # Table slot — right ~1/3
    tbl_left = img_left + img_w + Inches(0.15)
    tbl_top = img_top
    tbl_w = SLIDE_W - tbl_left - Inches(0.3)
    tbl_h = img_h
    _add_marker(slide, "table_slot", tbl_left, tbl_top, tbl_w, tbl_h, "Data table")


def _build_divider(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_brand_bar(slide, 0, Inches(0.12))
    _add_brand_bar(slide, SLIDE_H - Inches(0.12), Inches(0.12))

    eyebrow = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), SLIDE_W - Inches(1.2), Inches(0.4))
    _set_text(eyebrow.text_frame, "SECTION", size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    title = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), SLIDE_W - Inches(1.2), Inches(1.2))
    title.name = "title_slot"
    _set_text(title.text_frame, "{{layer_name}}", size=36, bold=True, color=INK, align=PP_ALIGN.CENTER)


def main() -> None:
    """Generate the cover + content + divider sample slides and write template.pptx."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_cover(prs)
    _build_content(prs)
    _build_divider(prs)

    out = Path(__file__).resolve().parents[1] / "src" / "services" / "ppt_assets" / "template.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
